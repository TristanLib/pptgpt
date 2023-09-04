import gradio as gr
import openai
import requests
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import os

from dotenv import load_dotenv, find_dotenv
_ = load_dotenv(find_dotenv())

openai.api_key  = os.getenv('OPENAI_API_KEY')

# Unsplash API的URL
# UNSPLASH_URL = "https://api.unsplash.com/photos/random?query={}&client_id=YOUR_UNSPLASH_API_KEY"

def get_completion(prompt, model="gpt-3.5-turbo-16k"):
    print(openai.api_key)
    messages = [{"role":"user", "content": prompt}]
    response = openai.ChatCompletion.create(
        model=model,
        messages=messages,
        temperature=0
    )
    return response.choices[0].message["content"]

prompt = f"""
Generate an outline of a PowerPoint based on the text delimited by triple backticks \ 
into a single sentence.
```{text}```
"""
def generate_ppt(prompt):
    # 使用OpenAI API获取提纲
    outline = get_completion(prompt)

    # 创建PPT
    prs = Presentation()
    for point in outline:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = point

        # 从Unsplash获取图片
        # response = requests.get(UNSPLASH_URL.format(point))
        # image_url = response.json()["urls"]["small"]
        # image = requests.get(image_url).content
        # image_stream = BytesIO(image)

        # # 添加图片到PPT
        # slide.shapes.add_picture(image_stream, Inches(1), Inches(1), width=Inches(4))

    # 保存PPT到文件
    ppt_file = "generated_ppt.pptx"
    prs.save(ppt_file)

    return ppt_file

# 创建gradio界面
interface = gr.Interface(
    fn=generate_ppt, 
    inputs=gr.Textbox(lines=5, placeholder="请输入你的提示语..."), 
    outputs=gr.File(),
    live=False
)

interface.launch()
